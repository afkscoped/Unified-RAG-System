"""
Document Processing Utilities
Handles PDF, DOCX, TXT, PPTX, XLSX extraction with chunking
"""

from typing import Dict, List, Optional
import io
from pathlib import Path
from loguru import logger

# Conditional imports
try:
    import PyPDF2
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False


class DocumentProcessor:
    """Multi-format document processor with intelligent chunking"""
    
    def __init__(self, chunk_size: int = 512, chunk_overlap: int = 50):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    
    def process_file(self, file, filename: Optional[str] = None) -> Dict:
        """
        Process uploaded file and extract text
        
        Returns dict with:
        - full_text: Complete document text
        - chunks: List of text chunks for indexing
        - metadata: Document metadata
        - type: File type
        """
        if filename is None:
            filename = getattr(file, 'name', 'unknown.txt')
        
        file_ext = Path(filename).suffix.lower()
        
        try:
            if file_ext == '.pdf':
                return self._process_pdf(file, filename)
            elif file_ext == '.docx':
                return self._process_docx(file, filename)
            elif file_ext == '.txt':
                return self._process_txt(file, filename)
            elif file_ext == '.pptx':
                return self._process_pptx(file, filename)
            elif file_ext in ['.xlsx', '.xls']:
                return self._process_xlsx(file, filename)
            else:
                # Try as plain text
                return self._process_txt(file, filename)
        except Exception as e:
            logger.error(f"Error processing {filename}: {e}")
            return {
                'full_text': f"Error processing file: {e}",
                'chunks': [],
                'num_chunks': 0,
                'type': 'ERROR',
                'filename': filename,
                'error': str(e)
            }
    
    def _process_pdf(self, file, filename: str) -> Dict:
        """Extract text from PDF"""
        full_text = ""
        num_pages = 0
        
        # Try pdfplumber first (better extraction)
        if HAS_PDFPLUMBER:
            try:
                file.seek(0)
                with pdfplumber.open(io.BytesIO(file.read())) as pdf:
                    num_pages = len(pdf.pages)
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            full_text += text + "\n\n"
            except Exception as e:
                logger.warning(f"pdfplumber failed, trying PyPDF2: {e}")
                full_text = ""
        
        # Fallback to PyPDF2
        if not full_text and HAS_PYPDF:
            try:
                file.seek(0)
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file.read()))
                num_pages = len(pdf_reader.pages)
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text:
                        full_text += text + "\n\n"
            except Exception as e:
                logger.error(f"PyPDF2 failed: {e}")
        
        if not full_text:
            full_text = "[Could not extract text from PDF]"
        
        chunks = self._create_chunks(full_text, filename)
        
        return {
            'full_text': full_text,
            'chunks': chunks,
            'num_chunks': len(chunks),
            'type': 'PDF',
            'num_pages': num_pages,
            'filename': filename
        }
    
    def _process_docx(self, file, filename: str) -> Dict:
        """Extract text from DOCX"""
        if not HAS_DOCX:
            return self._unsupported_format(filename, 'DOCX')
        
        try:
            file.seek(0)
            doc = Document(io.BytesIO(file.read()))
            full_text = "\n\n".join([para.text for para in doc.paragraphs if para.text])
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            full_text = f"[Error extracting DOCX: {e}]"
        
        chunks = self._create_chunks(full_text, filename)
        
        return {
            'full_text': full_text,
            'chunks': chunks,
            'num_chunks': len(chunks),
            'type': 'DOCX',
            'filename': filename
        }
    
    def _process_txt(self, file, filename: str) -> Dict:
        """Extract text from TXT"""
        try:
            file.seek(0)
            content = file.read()
            if isinstance(content, bytes):
                full_text = content.decode('utf-8', errors='ignore')
            else:
                full_text = content
        except Exception as e:
            logger.error(f"TXT extraction failed: {e}")
            full_text = f"[Error extracting TXT: {e}]"
        
        chunks = self._create_chunks(full_text, filename)
        
        return {
            'full_text': full_text,
            'chunks': chunks,
            'num_chunks': len(chunks),
            'type': 'TXT',
            'filename': filename
        }
    
    def _process_pptx(self, file, filename: str) -> Dict:
        """Extract text from PPTX"""
        if not HAS_PPTX:
            return self._unsupported_format(filename, 'PPTX')
        
        try:
            file.seek(0)
            prs = Presentation(io.BytesIO(file.read()))
            
            full_text = ""
            num_slides = len(prs.slides)
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        full_text += shape.text + "\n"
                full_text += "\n"
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            full_text = f"[Error extracting PPTX: {e}]"
            num_slides = 0
        
        chunks = self._create_chunks(full_text, filename)
        
        return {
            'full_text': full_text,
            'chunks': chunks,
            'num_chunks': len(chunks),
            'type': 'PPTX',
            'num_slides': num_slides,
            'filename': filename
        }
    
    def _process_xlsx(self, file, filename: str) -> Dict:
        """Extract text from XLSX"""
        if not HAS_XLSX:
            return self._unsupported_format(filename, 'XLSX')
        
        try:
            file.seek(0)
            wb = openpyxl.load_workbook(io.BytesIO(file.read()), data_only=True)
            
            full_text = ""
            for sheet in wb.worksheets:
                full_text += f"## Sheet: {sheet.title}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) for cell in row if cell is not None])
                    if row_text.strip():
                        full_text += row_text + "\n"
                full_text += "\n"
        except Exception as e:
            logger.error(f"XLSX extraction failed: {e}")
            full_text = f"[Error extracting XLSX: {e}]"
        
        chunks = self._create_chunks(full_text, filename)
        
        return {
            'full_text': full_text,
            'chunks': chunks,
            'num_chunks': len(chunks),
            'type': 'XLSX',
            'filename': filename
        }
    
    def _unsupported_format(self, filename: str, format_type: str) -> Dict:
        """Return error for unsupported format"""
        return {
            'full_text': f"[{format_type} support not available. Install required package.]",
            'chunks': [],
            'num_chunks': 0,
            'type': format_type,
            'filename': filename,
            'error': f"{format_type} library not installed"
        }
    
    def _create_chunks(self, text: str, source: str) -> List[Dict]:
        """Split text into overlapping chunks for indexing"""
        if not text or len(text.strip()) == 0:
            return []
        
        words = text.split()
        chunks = []
        
        for i in range(0, len(words), self.chunk_size - self.chunk_overlap):
            chunk_words = words[i:i + self.chunk_size]
            chunk_text = " ".join(chunk_words)
            
            if len(chunk_text.strip()) > 20:  # Minimum chunk size
                chunks.append({
                    'id': f"{source}_chunk_{len(chunks)}",
                    'content': chunk_text,
                    'metadata': {
                        'source': source,
                        'chunk_index': len(chunks),
                        'start_word': i,
                        'end_word': i + len(chunk_words)
                    }
                })
        
        return chunks
